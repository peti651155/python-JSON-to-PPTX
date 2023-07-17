import json
from tkinter import TOP
from pptx import Presentation
from pptx.util import Inches


def create_title_slide(slide, title, content):
    slide.shapes.title.text = title
    slide.placeholders[1].text = content


def create_text_slide(slide, title, content):
    slide.shapes.title.text = title
    slide.placeholders[1].text = content


def create_list_slide(slide, title, content):
    slide.shapes.title.text = title
    content_placeholder = slide.placeholders[1]
    for item in content:
        level = item['level']
        text = item['text']
        p = content_placeholder.text_frame.add_paragraph()
        p.text = text
        p.level = level


def create_picture_slide(slide, title, content):
    slide.shapes.title.text = title
    slide.shapes.add_picture(content, Inches(1), Inches(1), height=Inches(3))


def create_plot_slide(slide, title, content, configuration):
    slide.shapes.title.text = title
    slide.shapes.placeholders[0].text = f"X-label: {configuration['x-label']}"
    slide.shapes.placeholders[1].text = f"Y-label: {configuration['y-label']}"



def create_presentation(json_data): 
    presentation = Presentation()
    slides_data = json_data['presentation']

    for slide_data in slides_data:
        slide_type = slide_data['type']
        title = slide_data['title']
        content = slide_data['content']

        if slide_type == 'title':
            slide = presentation.slides.add_slide(presentation.slide_layouts[0])
            create_title_slide(slide, title, content)
        elif slide_type == 'text':
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            create_text_slide(slide, title, content)
        elif slide_type == 'list':
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            create_list_slide(slide, title, content)
        elif slide_type == 'picture':
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            create_picture_slide(slide, title, content)
        elif slide_type == 'plot':
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            configuration = slide_data['configuration']
            create_plot_slide(slide, title, content, configuration)

    return presentation


def main():
    with open('sample.json') as file:
        json_data = json.load(file)

    presentation = create_presentation(json_data)
    presentation.save('sample.pptx')

     
if __name__ == '__main__':
    main()
